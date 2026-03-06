

import boto3, pdfplumber, io, re, os, json, unicodedata, logging
from pathlib import Path

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

from dotenv import load_dotenv
load_dotenv()

AWS_ACCESS_KEY_ID     = os.environ.get("AWS_ACCESS_KEY_ID")
AWS_SECRET_ACCESS_KEY = os.environ.get("AWS_SECRET_ACCESS_KEY")
AWS_REGION            = os.environ.get("AWS_DEFAULT_REGION", "us-east-1")
BUCKET_NAME           = os.environ.get("BUCKET_NAME", "ocr-ai-for-bharat1")

_s3 = _textract = _bedrock = None
def _get_s3():
    global _s3
    if not _s3: _s3 = boto3.client("s3", region_name=AWS_REGION,
        aws_access_key_id=AWS_ACCESS_KEY_ID,
        aws_secret_access_key=AWS_SECRET_ACCESS_KEY)
    return _s3
def _get_textract():
    global _textract
    if not _textract: _textract = boto3.client("textract", region_name=AWS_REGION,
        aws_access_key_id=AWS_ACCESS_KEY_ID,
        aws_secret_access_key=AWS_SECRET_ACCESS_KEY)
    return _textract
def _get_bedrock():
    global _bedrock
    if not _bedrock: _bedrock = boto3.client("bedrock-runtime", region_name=AWS_REGION)
    return _bedrock

IS_LAMBDA       = bool(os.environ.get("AWS_LAMBDA_FUNCTION_NAME"))
import platform
if platform.system() == "Windows":
    TESSERACT_CMD = os.environ.get("TESSERACT_CMD", r"C:\Users\Hriday\Downloads\Release-25.12.0-0\poppler-25.12.0\Library\bin")
else:
    TESSERACT_CMD = os.environ.get("TESSERACT_CMD", "/usr/bin/tesseract")
USE_TEXTRACT    = os.environ.get("USE_TEXTRACT",  "true").lower()  == "true"
USE_TESSERACT   = os.environ.get("USE_TESSERACT", "true").lower()  == "true"
USE_HAIKU       = os.environ.get("USE_HAIKU",     "false").lower() == "true"
HAIKU_THRESHOLD = int(os.environ.get("HAIKU_THRESHOLD", "35"))
HAIKU_MAX_CHARS = int(os.environ.get("HAIKU_MAX_CHARS",  "2000"))
HAIKU_MODEL     = "anthropic.claude-3-haiku-20240307-v1:0"
IMAGE_EXTS      = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}
DOC_EXTS        = {".pdf", ".pptx", ".ppt", ".docx", ".doc"} | IMAGE_EXTS

# ── PDF ──────────────────────────────────────────────────────────────────────
def _pdf_digital(b):
    parts, found = [], 0
    try:
        with pdfplumber.open(io.BytesIO(b)) as pdf:
            for p in pdf.pages:
                t = p.extract_text()
                if t and len(t.strip()) > 20:
                    parts.append(t); found += 1
    except Exception as e: logger.warning(f"pdfplumber: {e}")
    return "\n".join(parts), found

def _pdf_tesseract(b):
    if not USE_TESSERACT or IS_LAMBDA: return ""
    try:
        import pytesseract; from pdf2image import convert_from_bytes
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        pages = []
        for i, img in enumerate(convert_from_bytes(b, dpi=200)):
            t = pytesseract.image_to_string(img, config="--oem 3 --psm 3 -l eng")
            if t.strip(): pages.append(t)
        return "\n".join(pages)
    except Exception as e: logger.warning(f"PDF tesseract: {e}"); return ""

# ── PPTX ─────────────────────────────────────────────────────────────────────
def _pptx_native(b):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(b))
        slides = []
        for i, slide in enumerate(prs.slides):
            parts = []
            if slide.shapes.title and slide.shapes.title.text.strip():
                parts.append(f"## {slide.shapes.title.text.strip()}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(r.text for r in para.runs).strip()
                        if line and line != getattr(slide.shapes.title, "text", "").strip():
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells: parts.append(" | ".join(cells))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes: parts.append(f"[Notes: {notes}]")
            if parts: slides.append(f"--- Slide {i+1} ---\n" + "\n".join(parts))
        return "\n\n".join(slides)
    except Exception as e: logger.warning(f"PPTX native: {e}"); return ""

# ── DOCX ─────────────────────────────────────────────────────────────────────
def _docx_native(b):
    try:
        from docx import Document
        doc = Document(io.BytesIO(b))
        parts = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if not t: continue
            lvl = para.style.name
            if lvl.startswith("Heading"):
                n = lvl.replace("Heading ", "")
                parts.append(f"{'#'*int(n) if n.isdigit() else '##'} {t}")
            else:
                parts.append(t)
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells: parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e: logger.warning(f"DOCX native: {e}"); return ""

# ── Image ────────────────────────────────────────────────────────────────────
def _image_tesseract(b):
    if not USE_TESSERACT: return ""
    try:
        import pytesseract; from PIL import Image
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        img = Image.open(io.BytesIO(b)).convert("RGB")
        return pytesseract.image_to_string(img, config="--oem 3 --psm 3 -l eng")
    except Exception as e: logger.warning(f"Image tesseract: {e}"); return ""

# ── Textract (shared fallback) ────────────────────────────────────────────────
def _textract_s3(bucket, key):
    if not USE_TEXTRACT: return ""
    logger.info("[OCR] Textract fallback ($0.0015/page)...")
    try:
        resp = _get_textract().detect_document_text(
            Document={"S3Object": {"Bucket": bucket, "Name": key}})
        return "\n".join(b["Text"] for b in resp["Blocks"] if b["BlockType"] == "LINE")
    except Exception as e: logger.error(f"Textract: {e}"); return ""

# ── Quality & cleanup ────────────────────────────────────────────────────────
def _quality_score(text):
    if not text or len(text) < 50: return 0
    words = text.split(); tw = max(len(words), 1); tc = max(len(text), 1)
    gc = sum(1 for c in text if ord(c)>127 and
             unicodedata.category(c) not in ("Ll","Lu","Lt","Lo","Nd"))
    nt = sum(1 for w in words if len(w)==1 and w not in ("a","I","-","•"))
    s  = 100 - min(40, int(gc/tc*400)) - min(30, int(nt/tw*300))
    s -= 10 if (sum(len(w) for w in words)/tw) < 3.0 else 0
    return max(0, min(100, s))

_FIXES = [(re.compile(p), r) for p, r in [
    (r"(?<=[a-zA-Z])0(?=[a-zA-Z])", "o"),
    (r"(?<=[a-zA-Z])1(?=[a-zA-Z])", "l"),
    (r"\bvv\b", "w"), (r"rn(?=[aeiou])", "m"),
    (r"(?<=\w)- (?=\w)", ""), (r" {2,}", " "),
]]
def _clean(text):
    out = []
    for line in text.split("\n"):
        s = line.strip()
        if not s: out.append(""); continue
        if re.fullmatch(r"\d{1,4}", s): continue
        if re.fullmatch(r"[\-\._ =|*~]{3,}", s): continue
        if len(s.split()) < 3 and not (s.isupper() or s.istitle() or s.endswith(":")): continue
        for p, r in _FIXES: s = p.sub(r, s)
        out.append(unicodedata.normalize("NFKC", s))
    return re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip()

def _haiku_cleanup(text):
    if not USE_HAIKU: return text
    if _quality_score(text) >= HAIKU_THRESHOLD: return text
    logger.info("[OCR] Haiku cleanup (1 call, capped)...")
    try:
        body = json.dumps({"anthropic_version":"bedrock-2023-05-31","max_tokens":512,
            "messages":[{"role":"user","content":
                f"Fix OCR errors. Keep all content. Output clean text ONLY.\n\n{text[:HAIKU_MAX_CHARS]}"}]})
        resp = _get_bedrock().invoke_model(modelId=HAIKU_MODEL,
            contentType="application/json", accept="application/json", body=body)
        cleaned = json.loads(resp["body"].read())["content"][0]["text"]
        return (cleaned + ("\n"+text[HAIKU_MAX_CHARS:] if len(text)>HAIKU_MAX_CHARS else "")).strip()
    except Exception as e: logger.warning(f"Haiku: {e}"); return text

# ── PUBLIC API ────────────────────────────────────────────────────────────────
def extract_text(bucket: str, key: str) -> str:
    """Called by main_handler.py: text = ocr.extract_text(BUCKET, f'uploads/{filename}')"""
    ext = Path(key).suffix.lower()
    if ext not in DOC_EXTS:
        raise ValueError(f"Unsupported file type: {ext}")

    logger.info(f"[OCR] START {ext} | s3://{bucket}/{key}")
    b = _get_s3().get_object(Bucket=bucket, Key=key)["Body"].read()

    if ext == ".pdf":
        text, found = _pdf_digital(b)
        if found == 0 or len(text.strip()) < 150:
            text = _pdf_tesseract(b)
        if len(text.strip()) < 150:
            text = _textract_s3(bucket, key)

    elif ext in (".pptx", ".ppt"):
        text = _pptx_native(b)
        if len(text.strip()) < 100:
            text = _textract_s3(bucket, key)

    elif ext in (".docx", ".doc"):
        text = _docx_native(b)
        if len(text.strip()) < 50:
            text = _textract_s3(bucket, key)

    elif ext in IMAGE_EXTS:
        text = _image_tesseract(b)
        if len(text.strip()) < 30:
            text = _textract_s3(bucket, key)

    if not text.strip():
        raise ValueError(f"No text found in s3://{bucket}/{key}")

    text = _clean(text)
    text = _haiku_cleanup(text)

    logger.info(f"[OCR] DONE: {len(text):,} chars | quality={_quality_score(text)}/100")
    return text